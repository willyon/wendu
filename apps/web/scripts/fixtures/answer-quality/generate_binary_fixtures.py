#!/usr/bin/env python3
"""生成问牍回答质量测试用的 pdf/docx/pptx/csv/xlsx（含可断言的植入事实）。"""
from pathlib import Path

OUT = Path(__file__).resolve().parent


def write_pdf():
    # 简易 ASCII PDF，便于 pypdf 抽取（避免未嵌入中文字体）
    line = "WENDU-QA-PDF-9182 owner=LinQi warehouse=North-Dock-7"
    # PDF string escape
    safe = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 50 720 Td ({safe}) Tj ET"
    stream_bytes = stream.encode("latin-1")
    objects = []
    objects.append(b"1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n")
    objects.append(b"2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n")
    objects.append(
        b"3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources<< /Font<< /F1 5 0 R >> >> >>endobj\n"
    )
    objects.append(
        f"4 0 obj<< /Length {len(stream_bytes)} >>stream\n".encode("latin-1")
        + stream_bytes
        + b"\nendstream\nendobj\n"
    )
    objects.append(b"5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n")

    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(out))
        out.extend(obj)
    xref_pos = len(out)
    out.extend(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    out.extend(b"0000000000 65535 f \n")
    for off in offsets[1:]:
        out.extend(f"{off:010d} 00000 n \n".encode("latin-1"))
    out.extend(
        f"trailer<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF\n".encode(
            "latin-1"
        )
    )
    (OUT / "planted-brief.pdf").write_bytes(bytes(out))


def write_docx():
    from docx import Document

    doc = Document()
    doc.add_heading("问牍质检 · Word 材料", level=1)
    doc.add_paragraph("唯一工单号：WENDU-DOCX-7741")
    doc.add_paragraph("审批人：周晚")
    doc.add_paragraph("结论：本材料仅用于回答质量回归，确认 Word 解析与检索可用。")
    doc.save(OUT / "planted-brief.docx")


def write_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(3))
    tf = box.text_frame
    tf.text = "演示密语 WENDU-PPTX-ALPHA-5520"
    p = tf.add_paragraph()
    p.text = "讲者：顾青禾 · 主题：多格式切片质检"
    p.font.size = Pt(18)
    prs.save(OUT / "planted-brief.pptx")


def write_csv():
    (OUT / "planted-inventory.csv").write_text(
        "sku,name,qty,note\n"
        "SKU-WENDU-CSV-3301,沧澜墨水,42,质检植入行\n"
        "SKU-OTHER-0001,普通纸,3,对照行\n",
        encoding="utf-8",
    )


def write_xlsx():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "库存"
    ws.append(["仓位", "物料", "数量"])
    ws.append(["A-17-WENDU-XLSX", "归档夹", 19])
    ws.append(["B-02", "普通订书机", 2])
    wb.save(OUT / "planted-inventory.xlsx")


if __name__ == "__main__":
    write_pdf()
    write_docx()
    write_pptx()
    write_csv()
    write_xlsx()
    print("generated:", ", ".join(sorted(p.name for p in OUT.glob("planted-*"))))
