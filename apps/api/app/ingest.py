"""Ingest: read file → extract text → chunk → local embed → write chunks.

Owns: per-format extract, chunking (Markdown / line-preserving tables), embed_passages,
files.status updates.
Does not own: upload auth / signed PUT (routers_files, storage), ask-time retrieve (rag).
Entry points: ingest_file(file_id) · chunk_text()
See: docs/implementation (ZH) §3–4.
"""

from io import BytesIO, StringIO
import csv
import re
from uuid import UUID

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import delete, text
from sqlalchemy.orm import Session

from .config import settings
from .embed import embed_passages
from .errors import AppError
from .models import Chunk, File
from .storage import get_storage


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_docx(data: bytes) -> str:
    doc = Document(BytesIO(data))
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_bits = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t:
                    slide_bits.append(t)
        if slide_bits:
            parts.append(f"[Slide {i}]")
            parts.extend(slide_bits)
    return "\n".join(parts)


def _extract_csv(data: bytes) -> str:
    raw = _decode_text(data)
    reader = csv.reader(StringIO(raw))
    rows = []
    for row in reader:
        cells = [str(c).strip() for c in row if str(c).strip()]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _extract_xlsx(data: bytes) -> str:
    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    parts = []
    try:
        for sheet in wb.worksheets:
            parts.append(f"[Sheet {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = []
                for cell in row:
                    if cell is None:
                        continue
                    s = str(cell).strip()
                    if s:
                        cells.append(s)
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def extract_text(filename: str, data: bytes) -> str:
    lower = filename.lower()
    if lower.endswith(".md") or lower.endswith(".markdown") or lower.endswith(".txt"):
        return _decode_text(data)
    if lower.endswith(".pdf"):
        reader = PdfReader(BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if lower.endswith(".docx"):
        return _extract_docx(data)
    if lower.endswith(".pptx"):
        return _extract_pptx(data)
    if lower.endswith(".csv"):
        return _extract_csv(data)
    if lower.endswith(".xlsx"):
        return _extract_xlsx(data)
    raise ValueError(f"unsupported type: {filename}")


def _split_long_text(text_value: str) -> list[str]:
    size = settings.chunk_size
    overlap = settings.chunk_overlap
    chunks = []
    i = 0
    n = len(text_value)
    while i < n:
        chunks.append(text_value[i : i + size])
        i += max(size - overlap, 1)
    return chunks


def _chunk_markdown(raw: str) -> list[str]:
    text = raw.replace("\r\n", "\n").replace("\r", "\n")
    if not text.strip():
        return []

    size = settings.chunk_size
    blocks: list[str] = []
    current: list[str] = []
    for line in text.split("\n"):
        if re.match(r"^#{1,6}\s", line) and current:
            blocks.append("\n".join(current).strip())
            current = [line]
        elif not line.strip() and current:
            blocks.append("\n".join(current).strip())
            current = []
        else:
            current.append(line)
    if current:
        blocks.append("\n".join(current).strip())

    blocks = [b for b in blocks if b]
    if not blocks:
        return []

    chunks: list[str] = []
    buf = ""
    for block in blocks:
        if len(block) <= size:
            if buf and len(buf) + 2 + len(block) <= size:
                buf = f"{buf}\n\n{block}"
            else:
                if buf:
                    chunks.append(buf)
                buf = block
        else:
            if buf:
                chunks.append(buf)
                buf = ""
            chunks.extend(_split_long_text(block))
    if buf:
        chunks.append(buf)
    return chunks


def _chunk_line_preserving(raw: str) -> list[str]:
    """Keep line breaks (tables, CSV, spreadsheet rows)."""
    text = raw.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        return []

    size = settings.chunk_size
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return []

    chunks: list[str] = []
    buf = ""
    for line in lines:
        if len(line) > size:
            if buf:
                chunks.append(buf)
                buf = ""
            chunks.extend(_split_long_text(line))
            continue
        if buf and len(buf) + 1 + len(line) > size:
            chunks.append(buf)
            buf = line
        else:
            buf = f"{buf}\n{line}" if buf else line
    if buf:
        chunks.append(buf)
    return chunks


def chunk_text(raw: str, filename: str = "") -> list[str]:
    lower = filename.lower()
    if lower.endswith(".md") or lower.endswith(".markdown"):
        pieces = _chunk_markdown(raw)
        if pieces:
            return pieces

    if lower.endswith(".csv") or lower.endswith(".xlsx"):
        pieces = _chunk_line_preserving(raw)
        if pieces:
            return pieces

    text_value = " ".join(raw.split())
    if not text_value:
        return []
    return _split_long_text(text_value)


def ingest_file(db: Session, file_id: UUID) -> None:
    storage = get_storage()
    f = db.get(File, file_id)
    if not f or f.status == "ready":
        return
    if f.status != "processing":
        f.status = "processing"
        f.fail_reason = None
        db.commit()
    try:
        data = storage.get_bytes(f.storage_key)
        raw = extract_text(f.filename, data)
        pieces = chunk_text(raw, f.filename)
        if not pieces:
            f.status = "failed"
            f.fail_reason = "PARSE_FAILED"
            db.commit()
            return
        vectors = embed_passages(pieces)
        db.execute(delete(Chunk).where(Chunk.file_id == f.id))
        for i, (piece, vec) in enumerate(zip(pieces, vectors)):
            db.add(
                Chunk(
                    user_id=f.user_id,
                    file_id=f.id,
                    ordinal=i,
                    text=piece,
                    embedding=vec,
                )
            )
        db.flush()
        db.execute(
            text(
                "UPDATE chunks SET tsv = to_tsvector('simple', text) WHERE file_id = :fid"
            ),
            {"fid": f.id},
        )
        f.status = "ready"
        f.fail_reason = None
        db.commit()
    except AppError as exc:
        db.rollback()
        f = db.get(File, file_id)
        if f:
            f.status = "failed"
            f.fail_reason = exc.code
            db.commit()
    except Exception:
        db.rollback()
        f = db.get(File, file_id)
        if f:
            f.status = "failed"
            f.fail_reason = "PARSE_FAILED"
            db.commit()
